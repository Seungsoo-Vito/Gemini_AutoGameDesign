import streamlit as st
import google.generativeai as genai
import requests
import base64
import json
import zlib
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import streamlit.components.v1 as components

# 1. 페이지 설정
st.set_page_config(page_title="비토쨩 GDD Pro", page_icon="🎮", layout="wide")

# --- 🎨 커스텀 CSS (부드러운 파스텔 & 프리미엄 가독성) ---
st.markdown("""
    <style>
    @import url('https://cdn.jsdelivr.net/gh/orioncactus/pretendard/dist/web/static/pretendard.css');
    
    .stApp {
        background-color: #fdfdfd;
        color: #2d3436;
        font-family: 'Pretendard', -apple-system, sans-serif;
    }

    [data-testid="stSidebar"] {
        background-color: #f1f3f5;
        border-right: 1px solid #e9ecef;
    }
    
    .main-title {
        font-size: 3.5rem !important;
        font-weight: 900 !important;
        background: linear-gradient(135deg, #ff9a9e 0%, #fad0c4 50%, #a1c4fd 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        letter-spacing: -0.05rem;
        margin-bottom: 0.5rem !important;
    }
    
    .gdd-card {
        background: #ffffff;
        border: 1px solid #f1f3f5;
        border-radius: 24px;
        padding: 40px;
        margin-bottom: 30px;
        box-shadow: 0 12px 24px rgba(149, 157, 165, 0.1);
    }
    
    .gdd-card h1, .gdd-card h2, .gdd-card h3 {
        color: #1f2937 !important;
        margin-top: 1.5rem !important;
        border-bottom: 3px solid #e0e7ff;
        display: inline-block;
        padding-bottom: 4px;
    }

    div.stButton > button[kind="primary"] {
        background: linear-gradient(135deg, #a1c4fd 0%, #c2e9fb 100%);
        color: #4b5563;
        border: none;
        border-radius: 14px;
        font-weight: 700;
        height: 3.5rem;
    }

    /* 구글 슬라이드 스타일 버튼 */
    .slide-export-btn {
        display: inline-flex;
        align-items: center;
        justify-content: center;
        background-color: #ffffff;
        color: #34a853;
        border: 2px solid #34a853;
        padding: 15px 30px;
        border-radius: 15px;
        text-decoration: none;
        font-weight: 800;
        font-size: 1.2rem;
        transition: all 0.3s;
        width: 100%;
        margin-top: 15px;
    }
    .slide-export-btn:hover {
        background-color: #34a853;
        color: white;
    }
    </style>
    """, unsafe_allow_html=True)

# --- 🔒 API 설정 ---
# 승수님, 여기에 API 키를 입력해 주세요.
API_KEY = "AIzaSyBpUR0gl_COhxbFPWxTiW6JJMuGgDF4Ams" 

if API_KEY.strip():
    genai.configure(api_key=API_KEY.strip())

# --- 📊 슬라이드 파일(PPTX) 생성 로직 ---
def create_presentation(slide_data, filename="GDD_Presentation.pptx"):
    prs = Presentation()
    
    # 타이틀 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Game Design Document"
    subtitle.text = "AI-Generated Strategic Pitch Deck"

    # AI가 생성한 슬라이드 데이터를 바탕으로 추가
    for item in slide_data:
        slide_layout = prs.slide_layouts[1] # 제목 및 본문 레이아웃
        slide = prs.slides.add_slide(slide_layout)
        
        # 제목 설정
        title_shape = slide.shapes.title
        title_shape.text = item.get("title", "Section")
        
        # 본문 설정 (불렛 포인트)
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        content = item.get("content", [])
        if isinstance(content, list):
            for i, line in enumerate(content):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = line
                p.level = 0
        else:
            tf.text = content

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- 🎨 이미지 생성 함수 (Imagen 4.0) ---
def generate_game_image(prompt_text):
    current_key = API_KEY.strip()
    if not current_key: return None
    url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-4.0-generate-001:predict?key={current_key}"
    payload = {"instances": [{"prompt": prompt_text}], "parameters": {"sampleCount": 1}}
    try:
        response = requests.post(url, json=payload, timeout=30)
        return response.json()["predictions"][0]["bytesBase64Encoded"]
    except: return None

# 세션 상태 관리
if 'gdd_result' not in st.session_state: st.session_state['gdd_result'] = None
if 'pptx_file' not in st.session_state: st.session_state['pptx_file'] = None
if 'generated_images' not in st.session_state: st.session_state['generated_images'] = {}

# --- 3. 메인 UI ---
st.markdown('<h1 class="main-title">비토쨩 GDD Pro 🎮</h1>', unsafe_allow_html=True)
st.write("기획서 작성부터 구글 슬라이드로 즉시 변환 가능한 발표 자료 생성까지")
st.divider()

# 입력 섹션 (4개 세션 완벽 유지)
with st.container():
    st.markdown('<div class="gdd-card">', unsafe_allow_html=True)
    c1, c2, c3, c4 = st.columns([1, 1, 1, 1])
    with c1:
        genre = st.selectbox("게임 장르", ["방치형 RPG", "수집형 RPG", "오픈월드", "로그라이크", "매치3 퍼즐", "액션 어드벤처"])
    with c2:
        target = st.selectbox("타겟 국가", ["글로벌", "한국", "일본", "북미/유럽"])
    with c3:
        art = st.selectbox("아트 스타일", ["픽셀 아트", "2D 카툰", "실사풍", "3D 캐주얼", "사이버펑크"])
    with c4:
        key = st.text_input("핵심 키워드", placeholder="예: 고양이, 타임루프")
    
    st.write("")
    if st.button("전문 기획서 및 슬라이드 생성 ✨", type="primary", use_container_width=True):
        if not API_KEY.strip():
            st.error("상단 API_KEY를 입력해 주세요.")
        elif not key:
            st.warning("키워드를 입력해 주세요.")
        else:
            with st.spinner("비토쨩이 기획서를 작성하고 슬라이드 문서를 생성 중입니다..."):
                model = genai.GenerativeModel('gemini-flash-latest')
                
                # 1. GDD 본문 생성
                gdd_prompt = f"당신은 시니어 게임 PM입니다. {genre}, {target}, {key}, {art} 조건으로 전문 GDD를 한국어로 작성하세요."
                gdd_res = model.generate_content(gdd_prompt)
                st.session_state['gdd_result'] = gdd_res.text
                
                # 2. 슬라이드 데이터 구조화 (JSON으로 요청)
                slide_prompt = f"""
                다음 기획서 내용을 바탕으로 발표용 슬라이드 10장 구성을 JSON 리스트 형식으로 만들어줘.
                형식: [{{"title": "슬라이드 제목", "content": ["불렛1", "불렛2"]}}, ...]
                기획서 내용: {gdd_res.text}
                """
                try:
                    # 구조화된 응답 유도 (MIME Type 사용 가능 시 설정)
                    slide_res = model.generate_content(slide_prompt)
                    # 텍스트에서 JSON 부분만 추출하는 간단한 클리닝
                    json_str = slide_res.text.strip().replace("```json", "").replace("```", "")
                    slide_data = json.loads(json_str)
                    
                    # 3. PPTX 파일 빌드
                    st.session_state['pptx_file'] = create_presentation(slide_data)
                except:
                    st.warning("슬라이드 문서 자동 구성 중 일부 오류가 있었으나 본문은 생성되었습니다.")
                
                # 4. 이미지 생성
                p_main = f"Game concept art, {genre}, {art}, {key}. Pastel color."
                img_b64 = generate_game_image(p_main)
                if img_b64: st.session_state['generated_images']["main"] = img_b64

# 결과 섹션
if st.session_state['gdd_result']:
    st.divider()
    
    # 기획서 본문 카드
    st.markdown('<div class="gdd-card">', unsafe_allow_html=True)
    st.subheader("📝 생성된 기획서 본문")
    if st.session_state['generated_images'].get("main"):
        img_cols = st.columns([1, 2, 1])
        img_cols[1].image(base64.b64decode(st.session_state['generated_images']["main"]), width=600)
    st.markdown(st.session_state['gdd_result'])
    st.markdown('</div>', unsafe_allow_html=True)

    # 🚀 구글 슬라이드 제작 완료 섹션
    st.subheader("📂 구글 슬라이드 문서 제작 완료")
    st.success("비토쨩이 기획서 내용을 바탕으로 10장의 발표 슬라이드 문서를 직접 만들었습니다!")
    
    if st.session_state['pptx_file']:
        st.download_button(
            label="✅ 생성된 구글 슬라이드(PPTX) 문서 다운로드",
            data=st.session_state['pptx_file'],
            file_name=f"GDD_{key}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        st.info("💡 위 파일을 다운로드한 후, 구글 드라이브(drive.google.com)에 '드래그 앤 드롭' 하면 즉시 구글 슬라이드로 변환되어 열립니다.")

st.caption("비토쨩 GDD Pro | Google Slides Document Engine | Powered by Google AI")